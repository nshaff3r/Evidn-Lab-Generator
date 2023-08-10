import collections
import collections.abc
import matplotlib.pyplot as plt
import io
import calendar
import csv
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from constructor import Lab, Building
from setup import add_labs
from sys import exit


prs = Presentation('template.pptx')
if len(prs.slides) != int(input("How many labs do you have? ")):
    exit("Your template presentation contains a different number of slides than labs.\n"
         "Please fix the template presentation to have the same number of slides as labs.")
building = add_labs()
if building is None:
    exit()

if len(prs.slides) != len(building.labs):
    exit("Your template presentation contains a different number of slides than labs.\n"
         "Please fix the template presentation to have the same number of slides as labs.")

timeperiod = input("Enter how you want the date to appear (ex. May 20–May 26, 2023): ")
bestenergy = [0, ""]
for lab in building.labs:
    metric = lab.energy_saved/lab.baseline_avg
    if metric > bestenergy[0]:
        bestenergy[0] = metric
        bestenergy[1] = f"{lab.name[0]}ab {lab.name[1:]}"
labofweek = bestenergy[1]

# For output csv file
fields = ["Lab", "Week Average", "Months Average", "Energy Saved", "Percent Saved"]
rows = []

for i, slide in enumerate(prs.slides):
    # Writes lab name
    building.labs[i].name = f"{building.labs[i].name[0]}ab {building.labs[i].name[1:]}"
    slide.shapes[1].text_frame.text = building.labs[i].name
    slide.shapes[1].text_frame.fit_text(font_family='title', font_file='Poppins-Regular.ttf', max_size=33, bold=True)
    slide.shapes[1].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Writes "Weekly Energy Report"
    subtitle = slide.shapes.add_textbox(Inches(0.44), Inches(1), Inches(3), Inches(1))
    subtitle.text_frame.text = "Weekly Energy Report"
    subtitle.text_frame.fit_text(font_family='subtitle', font_file='Poppins-Regular.ttf', max_size=14, bold=True)

    # Writes week dates
    dates = slide.shapes.add_textbox(Inches(0.44), Inches(1.31), Inches(2), Inches(0.5))
    dates.text_frame.text = timeperiod
    dates.text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=11, bold=False)

    # Changes text depending on if energy usage increased or decreased
    change = ""
    if (building.labs[i].energy_saved > 0):
        change = "Your energy reduction this week equates to:*"
        slide.shapes.add_picture("decrease.png", Inches(6.1), Inches(2.5), Inches(1.072), Inches(0.8))
    else:
        change = "Your increase in energy usage this week equates to:*"
        slide.shapes.add_picture("increase.png", Inches(6.1), Inches(2.5), Inches(1.072), Inches(0.8))

    slide.shapes[2].text_frame.text = change
    slide.shapes[2].text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=15, bold=True)
    slide.shapes[2].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Adds Fume Hood Alerts
    alerts = slide.shapes.add_textbox(Inches(6.01), Inches(1), Inches(0.5), Inches(0.6))
    alerts.text_frame.text = str(building.labs[i].alerts)
    alerts.text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=28, bold=True)

    # Adds energy equivalency in miles
    miles = slide.shapes.add_textbox(Inches(0.44), Inches(3.2), Inches(1.74), Inches(0.5))
    miles.text_frame.text = "{:,.1f}".format(building.labs[i].miles)
    miles.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    miles.text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=20, bold=True)
    miles.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Adds energy equivalency in phones charged
    phones = slide.shapes.add_textbox(Inches(2.39), Inches(3.2), Inches(2), Inches(0.5))
    phones.text_frame.text = "{:,.0f}".format(building.labs[i].phones)
    phones.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    phones.text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=20, bold=True)
    phones.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Adds energy equivalency in homes
    homes = slide.shapes.add_textbox(Inches(4.79), Inches(3.2), Inches(1.29), Inches(0.5))
    homes.text_frame.text = "{:,.3f}".format(building.labs[i].homes)
    homes.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    homes.text_frame.fit_text(font_family='dates', font_file='Poppins-Regular.ttf', max_size=20, bold=True)
    homes.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Lab of the week
    best_lab = slide.shapes.add_textbox(Inches(5.53), Inches(7.42), Inches(1.5), Inches(0.5))
    best_lab.text_frame.text = labofweek
    best_lab.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    best_lab.text_frame.fit_text(font_family='lab', font_file='Poppins-Regular.ttf', max_size=20, bold=True)
    best_lab.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xF0, 0x7C, 0x34)

    # Energy comparision bar chart
    bar_colors = []
    data = {}
    size = 0
    # If it's an ignored lab, there should be 2 instead of 3 bars
    if building.labs[i] not in building.ignored:
        data = {f"{building.shorthand}\nAverage": building.average, "Baseline\nAverage": building.labs[i].baseline_avg, "Your\nLab": building.labs[i].week_avg}
        bar_colors = ["orange", "grey", "pink"]
        size = 3
    else:
        data = {"Baseline\nAverage": building.labs[i].baseline_avg, "Your\nLab": building.labs[i].week_avg}
        bar_colors = ["grey", "pink"]
        size = 2
    labs = list(data.keys())
    values = list(data.values())
    plt.figure(figsize=(15, 5))
    plt.barh(labs, values, color=bar_colors)
    plt.margins(.3)
    for spine in plt.gca().spines.values():
        spine.set_visible(False)
    plt.xticks([])
    for j in range(size):
        plt.text(values[j], j, f"{values[j]: .2f} MTCO2", size=30)
    plt.yticks(size=20, fontweight='bold')
    plt.title("ENERGY USAGE THIS WEEK", size=35, fontweight='bold')
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    slide.shapes.add_picture(image_stream, Inches(0.3), Inches(4.6), Inches(6), Inches(2))
    
    # Months line graph
    plt.figure(figsize=(10, 5))
    data = building.labs[i].months_avg
    names = list(data.keys())
    baseline = {}
    for month in names:
        baseline[month] = building.labs[i].baseline_avg
    baselineValues = list(baseline.values())
    # # For when data starts in May
    # names = list(map(lambda x:  x + 4 if x <= 8 else x - 8, names))

    # Switches from month indexes to names
    names = list(map(lambda x: calendar.month_abbr[x], names))
    plt.plot(names, baselineValues, color="grey")

    # Current data (instead of baseline)
    values = list(data.values())
    values = list(map(lambda x: None if x == 0 else x, values))
    # # For when data starts in May
    # values = values[4:] + [None, None, None, None]
    plt.plot(names, values, color="orange")
    for spine in plt.gca().spines.values():
        spine.set_visible(False)
    plt.legend(['Baseline', "Current"], fontsize='x-large')
    plt.ylabel("MTons CO2", size=20)
    plt.xticks(size=20)
    plt.yticks(size=15)

    ax = plt.gca()
    ax.set_ylim(0.5 * building.labs[i].baseline_avg, 1.5 * building.labs[i].baseline_avg)
    # Hide every other month
    temp = ax.xaxis.get_ticklabels()
    temp = list(set(temp) - set(temp[::2]))
    for label in temp:
        label.set_visible(False)
    plt.title("CURRENT VS. BASELINE ENERGY USAGE", size=27, fontweight='bold', pad=20)
    image_stream = io.BytesIO()
    plt.savefig(image_stream)
    slide.shapes.add_picture(image_stream, Inches(0.3), Inches(6.8), Inches(5.13), Inches(2.58))
    rows.append([building.labs[i].name, building.labs[i].week_avg, dict(zip(names, values)),
                 building.labs[i].energy_saved,
                 building.labs[i].energy_saved/building.labs[i].baseline_avg * 100])

prs.save("output.pptx")
with open("data.csv", "w") as csvfile:
    csvwriter = csv.writer(csvfile)
    csvwriter.writerow(fields) 
    csvwriter.writerows(rows)
print("Finished! Generated output.pptx and data.csv.")